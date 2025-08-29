from __future__ import annotations

import argparse
from datetime import datetime
from pathlib import Path
import yaml
from pptx import Presentation
from pptx.util import Inches


def main(argv=None):
    p = argparse.ArgumentParser()
    p.add_argument("--outline", required=True)
    p.add_argument("--shots", required=True)
    p.add_argument("--out", required=True)
    args = p.parse_args(argv)

    outline = yaml.safe_load(Path(args.outline).read_text())
    prs = Presentation()
    for slide in outline:
        layout = prs.slide_layouts[1] if slide.get("bullets") else prs.slide_layouts[0]
        s = prs.slides.add_slide(layout)
        s.shapes.title.text = slide.get("title", "")
        if slide.get("bullets"):
            tf = s.shapes.placeholders[1].text_frame
            first = True
            for bullet in slide.get("bullets", []):
                if first:
                    tf.text = bullet
                    first = False
                else:
                    tf.add_paragraph().text = bullet
    for img in sorted(Path(args.shots).glob("*.png")):
        s = prs.slides.add_slide(prs.slide_layouts[5])
        s.shapes.add_picture(str(img), Inches(0), Inches(0), width=prs.slide_width)
    out_dir = Path(args.out)
    out_dir.mkdir(parents=True, exist_ok=True)
    stamp = datetime.utcnow().strftime("%Y%m%d_%H%M%S")
    pptx_path = out_dir / f"DR-RD_Walkthrough_{stamp}.pptx"
    prs.save(pptx_path)
    try:
        from reportlab.pdfgen import canvas

        pdf_path = out_dir / f"DR-RD_Walkthrough_{stamp}.pdf"
        c = canvas.Canvas(str(pdf_path))
        c.drawString(100, 750, "See PPTX for slides")
        c.save()
    except Exception:
        pass
    print(f"Deck written to {pptx_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
