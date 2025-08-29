import tempfile
from pathlib import Path
from pptx import Presentation

from scripts.demos import flows
from scripts import snapshots, generate_deck


def test_deck_built():
    run_dir = Path(tempfile.mkdtemp())
    flows.run_all(run_dir)
    shots_dir = Path(tempfile.mkdtemp())
    snapshots.main(["--runs", str(run_dir), "--out", str(shots_dir)])
    out_dir = Path(tempfile.mkdtemp())
    generate_deck.main([
        "--outline",
        "docs/templates/deck_outline.yaml",
        "--shots",
        str(shots_dir),
        "--out",
        str(out_dir),
    ])
    pptx_file = next(out_dir.glob("*.pptx"))
    prs = Presentation(str(pptx_file))
    titles = [s.shapes.title.text for s in prs.slides if s.shapes.title]
    assert "Demo Highlights" in titles
